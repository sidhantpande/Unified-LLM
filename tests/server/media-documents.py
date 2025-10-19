"""
Server Media Tests - Document Processing
Tests OpenAI-compatible endpoints with document files (PDF, Office docs).
"""

import pytest
import json
import tempfile
import requests
import os
from typing import Dict, Any, List
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# Test configuration
SERVER_BASE_URL = "http://localhost:8000"
TIMEOUT = 90  # Longer timeout for document processing

# Test models - using text models since document processing works across all models
TEXT_MODELS = {
    "ollama": [
        "qwen3:4b-instruct",
        "llama3:8b",
        "gemma3:2b"
    ],
    "lmstudio": [
        "qwen/qwen3-next-80b",
        "meta-llama/llama-3.2-8b-instruct",
        "microsoft/DialoGPT-large"
    ]
}

class DocumentTestHelper:
    """Helper class for document testing utilities."""

    @staticmethod
    def create_test_pdf(content: str = None, filename: str = None) -> str:
        """Create a test PDF file."""
        if content is None:
            content = """
            Test Document Title

            This is a test PDF document created for testing AbstractCore's media processing capabilities.

            Key Information:
            • Product: AbstractCore Media System
            • Version: 2.0
            • Features: Universal file processing, vision capabilities, streaming support
            • Status: Production Ready

            Technical Details:
            The system supports 12+ file formats including images, PDFs, Office documents,
            and data files with intelligent processing and graceful fallback handling.

            Conclusion:
            This document demonstrates successful PDF processing through the OpenAI-compatible
            server endpoints with automatic content extraction and analysis.
            """

        # Create temporary PDF
        temp_file = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        temp_file.close()

        if filename:
            pdf_path = filename
        else:
            pdf_path = temp_file.name

        # Create PDF with reportlab
        c = canvas.Canvas(pdf_path, pagesize=letter)
        width, height = letter

        # Add content
        lines = content.strip().split('\n')
        y_position = height - 50

        for line in lines:
            line = line.strip()
            if line:
                if line.isupper() or "Title" in line:
                    c.setFont("Helvetica-Bold", 14)
                else:
                    c.setFont("Helvetica", 10)

                c.drawString(50, y_position, line[:90])  # Limit line length
                y_position -= 15

                if y_position < 50:  # New page if needed
                    c.showPage()
                    y_position = height - 50

        c.save()
        return pdf_path

    @staticmethod
    def create_test_docx(content: str = None) -> str:
        """Create a test Word document."""
        if content is None:
            content = {
                "title": "AbstractCore Media Testing Document",
                "sections": [
                    {
                        "heading": "Overview",
                        "text": "This Word document tests the media processing capabilities of AbstractCore's server endpoints."
                    },
                    {
                        "heading": "Features Tested",
                        "text": "Document extraction, content analysis, and OpenAI-compatible API integration."
                    },
                    {
                        "heading": "Expected Results",
                        "text": "The system should extract and analyze this content automatically."
                    }
                ]
            }

        # Create Word document
        doc = Document()
        doc.add_heading(content["title"], 0)

        for section in content["sections"]:
            doc.add_heading(section["heading"], level=1)
            doc.add_paragraph(section["text"])

        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
        temp_file.close()
        doc.save(temp_file.name)

        return temp_file.name

    @staticmethod
    def create_test_xlsx() -> str:
        """Create a test Excel spreadsheet."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Sales Data"

        # Headers
        ws['A1'] = "Quarter"
        ws['B1'] = "Revenue"
        ws['C1'] = "Growth"
        ws['D1'] = "Region"

        # Data
        data = [
            ["Q1 2024", 150000, "12%", "North"],
            ["Q2 2024", 175000, "15%", "South"],
            ["Q3 2024", 200000, "18%", "East"],
            ["Q4 2024", 225000, "20%", "West"]
        ]

        for row_idx, row_data in enumerate(data, start=2):
            for col_idx, value in enumerate(row_data, start=1):
                ws.cell(row=row_idx, column=col_idx, value=value)

        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        temp_file.close()
        wb.save(temp_file.name)

        return temp_file.name

    @staticmethod
    def create_test_pptx() -> str:
        """Create a test PowerPoint presentation."""
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "AbstractCore Media Testing"
        subtitle.text = "Server Endpoint Document Processing Validation"

        # Content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Test Objectives"

        tf = body_shape.text_frame
        tf.text = "Validate document processing capabilities"

        p = tf.add_paragraph()
        p.text = "Test PowerPoint content extraction"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "Verify OpenAI API compatibility"
        p.level = 1

        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        temp_file.close()
        prs.save(temp_file.name)

        return temp_file.name

    @staticmethod
    def is_server_running() -> bool:
        """Check if the server is running."""
        try:
            response = requests.get(f"{SERVER_BASE_URL}/health", timeout=5)
            return response.status_code == 200
        except:
            return False

    @staticmethod
    def get_available_models(provider: str) -> List[str]:
        """Get available models for a provider from the server."""
        try:
            response = requests.get(f"{SERVER_BASE_URL}/providers/{provider}/models", timeout=10)
            if response.status_code == 200:
                data = response.json()
                return data.get("models", [])
            return []
        except:
            return []

@pytest.fixture(scope="module")
def server_check():
    """Ensure server is running before tests."""
    if not DocumentTestHelper.is_server_running():
        pytest.skip("Server is not running. Start with: uvicorn abstractcore.server.app:app --port 8000")

@pytest.fixture(scope="module")
def test_documents():
    """Create test documents for the test session."""
    documents = {}

    # Create test documents
    documents["pdf"] = DocumentTestHelper.create_test_pdf()
    documents["docx"] = DocumentTestHelper.create_test_docx()
    documents["xlsx"] = DocumentTestHelper.create_test_xlsx()
    documents["pptx"] = DocumentTestHelper.create_test_pptx()

    yield documents

    # Cleanup
    for doc_path in documents.values():
        try:
            if os.path.exists(doc_path):
                os.unlink(doc_path)
        except:
            pass

class TestDocumentProcessingOpenAI:
    """Test document processing using OpenAI format."""

    @pytest.mark.parametrize("provider", ["ollama", "lmstudio"])
    @pytest.mark.parametrize("doc_type", ["pdf", "docx", "xlsx", "pptx"])
    def test_document_analysis_openai_format(self, server_check, test_documents, provider, doc_type):
        """Test document analysis using OpenAI format with base64 encoding."""
        available_models = DocumentTestHelper.get_available_models(provider)
        text_models = [m for m in TEXT_MODELS.get(provider, []) if m in available_models]

        if not text_models:
            pytest.skip(f"No text models available for {provider}")

        model = text_models[0]
        doc_path = test_documents[doc_type]

        # Read document and encode as base64 (simulating file upload)
        with open(doc_path, "rb") as f:
            doc_data = f.read()

        # Note: For documents, we'll use the @filename syntax since OpenAI doesn't
        # natively support document types in image_url format
        payload = {
            "model": f"{provider}/{model}",
            "messages": [
                {
                    "role": "user",
                    "content": f"Please analyze the content of @{doc_path} and provide a brief summary."
                }
            ],
            "max_tokens": 300,
            "temperature": 0.1
        }

        response = requests.post(f"{SERVER_BASE_URL}/v1/chat/completions", json=payload, timeout=TIMEOUT)

        assert response.status_code == 200, f"Request failed for {doc_type}: {response.text}"

        data = response.json()
        assert "choices" in data
        assert len(data["choices"]) > 0
        assert "message" in data["choices"][0]
        assert "content" in data["choices"][0]["message"]

        content = data["choices"][0]["message"]["content"].lower()

        # Verify document content was processed
        expected_terms = {
            "pdf": ["test", "document", "abstractcore", "media", "processing"],
            "docx": ["abstractcore", "media", "testing", "document", "features"],
            "xlsx": ["sales", "data", "quarter", "revenue", "growth"],
            "pptx": ["abstractcore", "media", "testing", "objectives", "processing"]
        }

        doc_terms = expected_terms[doc_type]
        found_terms = [term for term in doc_terms if term in content]

        assert len(found_terms) >= 2, \
            f"Should find at least 2 expected terms in {doc_type} content. " \
            f"Expected: {doc_terms}, Found: {found_terms}, Content: {content[:200]}"

    @pytest.mark.parametrize("provider", ["ollama", "lmstudio"])
    def test_pdf_specific_analysis(self, server_check, test_documents, provider):
        """Test specific PDF analysis capabilities."""
        available_models = DocumentTestHelper.get_available_models(provider)
        text_models = [m for m in TEXT_MODELS.get(provider, []) if m in available_models]

        if not text_models:
            pytest.skip(f"No text models available for {provider}")

        model = text_models[0]

        payload = {
            "model": f"{provider}/{model}",
            "messages": [
                {
                    "role": "user",
                    "content": f"What is the main product mentioned in @{test_documents['pdf']}? List the key features."
                }
            ],
            "max_tokens": 200
        }

        response = requests.post(f"{SERVER_BASE_URL}/v1/chat/completions", json=payload, timeout=TIMEOUT)

        assert response.status_code == 200
        data = response.json()
        content = data["choices"][0]["message"]["content"].lower()

        # Should identify AbstractCore and its features
        assert any(term in content for term in ["abstractcore", "media", "system"]), \
            f"Should identify the main product. Got: {content}"

    @pytest.mark.parametrize("provider", ["ollama", "lmstudio"])
    def test_excel_data_analysis(self, server_check, test_documents, provider):
        """Test Excel spreadsheet data analysis."""
        available_models = DocumentTestHelper.get_available_models(provider)
        text_models = [m for m in TEXT_MODELS.get(provider, []) if m in available_models]

        if not text_models:
            pytest.skip(f"No text models available for {provider}")

        model = text_models[0]

        payload = {
            "model": f"{provider}/{model}",
            "messages": [
                {
                    "role": "user",
                    "content": f"Analyze the data in @{test_documents['xlsx']}. What trends do you see in the revenue growth?"
                }
            ],
            "max_tokens": 200
        }

        response = requests.post(f"{SERVER_BASE_URL}/v1/chat/completions", json=payload, timeout=TIMEOUT)

        assert response.status_code == 200
        data = response.json()
        content = data["choices"][0]["message"]["content"].lower()

        # Should identify revenue/growth trends
        assert any(term in content for term in ["revenue", "growth", "quarter", "increase", "trend"]), \
            f"Should identify revenue trends. Got: {content}"

class TestDocumentStreamingOpenAI:
    """Test streaming document processing."""

    @pytest.mark.parametrize("provider", ["ollama", "lmstudio"])
    def test_streaming_pdf_analysis(self, server_check, test_documents, provider):
        """Test streaming analysis of PDF documents."""
        available_models = DocumentTestHelper.get_available_models(provider)
        text_models = [m for m in TEXT_MODELS.get(provider, []) if m in available_models]

        if not text_models:
            pytest.skip(f"No text models available for {provider}")

        model = text_models[0]

        payload = {
            "model": f"{provider}/{model}",
            "messages": [
                {
                    "role": "user",
                    "content": f"Provide a detailed analysis of @{test_documents['pdf']}"
                }
            ],
            "stream": True,
            "max_tokens": 300
        }

        response = requests.post(f"{SERVER_BASE_URL}/v1/chat/completions", json=payload, stream=True, timeout=TIMEOUT)

        assert response.status_code == 200

        # Collect streaming response
        content_parts = []
        for line in response.iter_lines():
            if line:
                line_str = line.decode('utf-8')
                if line_str.startswith('data: ') and not line_str.endswith('[DONE]'):
                    try:
                        data = json.loads(line_str[6:])
                        if 'choices' in data and data['choices']:
                            delta = data['choices'][0].get('delta', {})
                            if 'content' in delta:
                                content_parts.append(delta['content'])
                    except json.JSONDecodeError:
                        pass

        full_content = ''.join(content_parts).lower()

        assert len(full_content) > 50, "Should receive substantial streaming content for document analysis"
        assert any(term in full_content for term in ["document", "abstractcore", "media", "processing"]), \
            f"Streaming content should reference document content. Got: {full_content[:200]}"

class TestDocumentErrorHandling:
    """Test error handling for document processing."""

    def test_missing_file_error(self, server_check):
        """Test error handling for missing files."""
        payload = {
            "model": "ollama/qwen3:4b-instruct",
            "messages": [
                {
                    "role": "user",
                    "content": "Analyze @/nonexistent/file.pdf"
                }
            ]
        }

        response = requests.post(f"{SERVER_BASE_URL}/v1/chat/completions", json=payload, timeout=TIMEOUT)

        assert response.status_code == 400
        data = response.json()
        assert "error" in data
        assert data["error"]["type"] == "file_not_found"

    def test_unsupported_file_type(self, server_check):
        """Test error handling for unsupported file types."""
        # Create a file with unsupported extension
        with tempfile.NamedTemporaryFile(suffix=".exe", delete=False) as f:
            f.write(b"fake executable content")
            temp_file = f.name

        try:
            payload = {
                "model": "ollama/qwen3:4b-instruct",
                "messages": [
                    {
                        "role": "user",
                        "content": f"Analyze @{temp_file}"
                    }
                ]
            }

            response = requests.post(f"{SERVER_BASE_URL}/v1/chat/completions", json=payload, timeout=TIMEOUT)

            assert response.status_code == 400
            data = response.json()
            assert "error" in data
            assert data["error"]["type"] == "invalid_file_type"

        finally:
            if os.path.exists(temp_file):
                os.unlink(temp_file)

    def test_file_too_large_error(self, server_check):
        """Test error handling for files that are too large."""
        # Create a large PDF (>10MB)
        large_content = "Large file test content. " * 500000  # ~12MB of text

        temp_file = tempfile.NamedTemporaryFile(suffix=".txt", delete=False)
        with open(temp_file.name, 'w') as f:
            f.write(large_content)

        try:
            payload = {
                "model": "ollama/qwen3:4b-instruct",
                "messages": [
                    {
                        "role": "user",
                        "content": f"Summarize @{temp_file.name}"
                    }
                ]
            }

            response = requests.post(f"{SERVER_BASE_URL}/v1/chat/completions", json=payload, timeout=TIMEOUT)

            # Should handle large files (might succeed with chunking or fail gracefully)
            assert response.status_code in [200, 400]

            if response.status_code == 400:
                data = response.json()
                assert "error" in data
                assert data["error"]["type"] in ["file_too_large", "total_size_exceeded"]

        finally:
            if os.path.exists(temp_file.name):
                os.unlink(temp_file.name)

if __name__ == "__main__":
    # Quick validation run
    helper = DocumentTestHelper()

    if not helper.is_server_running():
        print("❌ Server not running. Start with: uvicorn abstractcore.server.app:app --port 8000")
        exit(1)

    print("✅ Server is running")

    # Check available models
    for provider in ["ollama", "lmstudio"]:
        models = helper.get_available_models(provider)
        text_models = [m for m in TEXT_MODELS.get(provider, []) if m in models]
        print(f"📝 {provider.title()} text models available: {text_models}")

    # Test document creation
    try:
        print("\n🧪 Testing document creation...")
        pdf_path = helper.create_test_pdf()
        print(f"✅ Created test PDF: {pdf_path}")

        docx_path = helper.create_test_docx()
        print(f"✅ Created test DOCX: {docx_path}")

        xlsx_path = helper.create_test_xlsx()
        print(f"✅ Created test XLSX: {xlsx_path}")

        pptx_path = helper.create_test_pptx()
        print(f"✅ Created test PPTX: {pptx_path}")

        # Cleanup
        for path in [pdf_path, docx_path, xlsx_path, pptx_path]:
            if os.path.exists(path):
                os.unlink(path)

        print("✅ Document creation tests passed")

    except Exception as e:
        print(f"❌ Document creation failed: {e}")

    print("\n🧪 Run tests with: pytest tests/server/media-documents.py -v")